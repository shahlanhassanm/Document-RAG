"""
Optimized Multi-Format Document Processor.
Supports: PDF, DOCX, PPTX, XLSX, CSV, TXT
Implements structural data recovery, deterministic resource management, 
and resilient encoding handling.
"""

import os
import csv
import datetime
import logging
from typing import List, Tuple, Optional, Any, Iterator
from dataclasses import dataclass, field

from langchain_core.documents import Document
from core.logger import get_logger

log = get_logger("document_processor")

@dataclass
class FailedFile:
    path: str
    error_type: str
    message: str

@dataclass
class BatchResult:
    documents: List = field(default_factory=list)
    failures: List[FailedFile] = field(default_factory=list)
    stats: dict = field(default_factory=lambda: {
        "start_time": datetime.datetime.now(),
        "total_files": 0,
        "processed_files": 0
    })

def process_file(file_path: str, original_name: str = None) -> List:
    """
    Dispatcher for file processing.
    """
    ext = os.path.splitext(file_path).[1]lower()
    name = original_name or os.path.basename(file_path)
    
    # MIME-type/Magic byte check would go here for production hardening
    
    try:
        if ext == '.pdf':
            return _process_pdf(file_path, name)
        elif ext == '.docx':
            return _process_docx(file_path, name)
        elif ext == '.pptx':
            return _process_pptx(file_path, name)
        elif ext in ('.xlsx', '.xls'):
            return _process_excel(file_path, name)
        elif ext == '.csv':
            return _process_csv(file_path, name)
        elif ext in ('.txt', '.md'):
            return _process_txt(file_path, name)
        else:
            raise ValueError(f"Unsupported file format: {ext}")
    except Exception as e:
        log.error(f"Failed to process {name}: {str(e)}", exc_info=True)
        raise

# --- PDF PROCESSING ---

def _process_pdf(file_path: str, name: str) -> List:
    """Resilient PDF extraction using context manager."""
    import pymupdf  # Modern PyMuPDF
    docs =
    
    with pymupdf.open(file_path) as pdf:
        total_pages = len(pdf)
        for page_num, page in enumerate(pdf, 1):
            text = page.get_text().strip()
            
            # Placeholder for OCR fallback if text is empty
            if not text:
                log.debug(f"Page {page_num} of {name} is empty (Potential Image/Scan)")
            
            docs.append(Document(
                page_content=text,
                metadata={
                    "source": name,
                    "file_type": "pdf",
                    "page": page_num,
                    "total_segments": total_pages,
                    "extraction_timestamp": datetime.datetime.now().isoformat()
                }
            ))
    return docs

# --- WORD PROCESSING ---

def _process_docx(file_path: str, name: str) -> List:
    """Structure-aware DOCX extraction including tables, headers, and footers."""
    from docx import Document as DocxReader
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    
    doc = DocxReader(file_path)
    docs =
    
    # 1. Extract Headers/Footers from all sections
    header_footer_text =
    for section in doc.sections:
        if section.header and section.header.paragraphs:
            header_footer_text.extend([p.text for p in section.header.paragraphs if p.text.strip()])
        if section.footer and section.footer.paragraphs:
            header_footer_text.extend([p.text for p in section.footer.paragraphs if p.text.strip()])
    
    metadata_base = {
        "source": name,
        "file_type": "docx",
        "extraction_timestamp": datetime.datetime.now().isoformat()
    }

    # Helper to iterate through paragraphs and tables in document order
    def iter_block_items(parent):
        from docx.document import Document as _Document
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import Table, _Cell
        
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    full_content =
    if header_footer_text:
        full_content.append("Headers/Footers:\n" + "\n".join(header_footer_text))

    for block in iter_block_items(doc):
        if isinstance(block, Paragraph):
            if block.text.strip():
                full_content.append(block.text)
        elif isinstance(block, Table):
            table_data =
            for row in block.rows:
                row_text = " | ".join(cell.text.strip().replace('\n', ' ') for cell in row.cells)
                table_data.append(row_text)
            full_content.append("\n".join(table_data))

    # Chunking logic would typically be applied here. 
    # For now, we return as a single document with recovered structure.
    docs.append(Document(
        page_content="\n\n".join(full_content),
        metadata=metadata_base
    ))
    return docs

# --- POWERPOINT PROCESSING ---

def _process_pptx(file_path: str, name: str) -> List:
    """Structure-aware PPTX extraction including Speaker Notes and GroupShapes."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(file_path)
    docs =
    total_slides = len(prs.slides)

    def extract_text_recursive(shape) -> str:
        texts =
        # Standard Text Boxes
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text.strip())
        
        # Tables
        elif shape.has_table:
            for row in shape.table.rows:
                row_text = " | ".join(cell.text_frame.text.strip().replace('\n', ' ') for cell in row.cells)
                texts.append(row_text)
        
        # Recursive GroupShapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                texts.append(extract_text_recursive(sub_shape))
                
        return "\n".join(filter(None, texts))

    for i, slide in enumerate(prs.slides, 1):
        slide_content =
        
        # 1. Slide Shapes (Ordered)
        # Advanced: sort by y, then x coordinates
        for shape in slide.shapes:
            text = extract_text_recursive(shape)
            if text:
                slide_content.append(text)
        
        # 2. Speaker Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f": {notes}")

        if slide_content:
            docs.append(Document(
                page_content="\n".join(slide_content),
                metadata={
                    "source": name,
                    "file_type": "pptx",
                    "slide": i,
                    "total_segments": total_slides,
                    "extraction_timestamp": datetime.datetime.now().isoformat()
                }
            ))
    return docs

# --- EXCEL PROCESSING ---

def _process_excel(file_path: str, name: str) -> List:
    """Safe Excel extraction using context manager and formula result handling."""
    from openpyxl import load_workbook
    docs =
    
    with load_workbook(file_path, read_only=True, data_only=True) as wb:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows =
            for row in ws.iter_rows(values_only=True):
                # Clean newlines and handle None
                row_text = " | ".join([
                    str(cell).strip().replace('\n', ' ') if cell is not None else "" 
                    for cell in row
                ])
                # Skip rows that only contain separators
                if row_text.replace("|", "").strip():
                    rows.append(row_text)
            
            if rows:
                docs.append(Document(
                    page_content=f"Sheet: {sheet_name}\n" + "\n".join(rows),
                    metadata={
                        "source": name,
                        "file_type": "excel",
                        "sheet": sheet_name,
                        "extraction_timestamp": datetime.datetime.now().isoformat()
                    }
                ))
    return docs

# --- FLAT FILE PROCESSING ---

def _process_csv(file_path: str, name: str) -> List:
    """BOM-safe and newline-safe CSV extraction."""
    rows =
    # utf-8-sig handles potential BOM from Excel exports
    with open(file_path, 'r', encoding='utf-8-sig', newline='', errors='replace') as f:
        reader = csv.reader(f)
        for row in reader:
            row_text = " | ".join(row)
            if row_text.strip():
                rows.append(row_text)
                
    if not rows: return
    
    return

def _process_txt(file_path: str, name: str) -> List:
    """BOM-safe plain text extraction."""
    with open(file_path, 'r', encoding='utf-8-sig', errors='replace') as f:
        text = f.read().strip()
    
    if not text: return
    
    return

# --- BATCH RUNNER ---

def process_batch(file_paths: List]) -> BatchResult:
    """
    Processes a batch of files and returns a structured result for observability.
    Args: List of (file_path, original_name)
    """
    result = BatchResult()
    result.stats["total_files"] = len(file_paths)
    
    log.info(f"Starting batch: {len(file_paths)} files")
    
    for file_path, original_name in file_paths:
        try:
            docs = process_file(file_path, original_name)
            result.documents.extend(docs)
            result.stats["processed_files"] += 1
        except Exception as e:
            result.failures.append(FailedFile(
                path=original_name,
                error_type=type(e).__name__,
                message=str(e)
            ))
            log.error(f"Skipping {original_name} due to error.")

    duration = datetime.datetime.now() - result.stats["start_time"]
    log.info(f"Batch completed in {duration.total_seconds():.2f}s. "
             f"Success: {result.stats['processed_files']}, Failures: {len(result.failures)}")
    
    return result